from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Configuration
BLUE_COLOR = RGBColor(0x41, 0x5f, 0xff)  # #415fff
FONT_CN = "HarmonyOS Sans SC"
FONT_EN = "Inter"

def set_font(run, size, is_bold=False):
    run.font.size = size
    run.font.bold = is_bold
    run.font.name = FONT_EN
    rPr = run._r.get_or_add_rPr()
    ea = qn('a:ea')
    rPr.set(ea, FONT_CN)

def add_content_slide(prs, title_text, content_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # 1. Blue Title Box (Top Left)
    # Position and size
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(4.5)
    height = Inches(0.9)
    
    shape = slide.shapes.add_shape(
        1, # msoShapeRectangle
        left, top, width, height
    )
    
    # Fill color #415fff
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_COLOR
    shape.line.fill.background() # No outline
    
    # Title Text
    text_frame = shape.text_frame
    text_frame.vertical_anchor = 3 # Middle
    text_frame.margin_left = Inches(0.2)
    
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.color.rgb = RGBColor(255, 255, 255) # White text
    set_font(run, Pt(20), True)

    # 2. Main Content
    # Positioned below the header
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for line in content_text.strip().split('\n'):
        line = line.strip()
        if not line: continue
        
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        p.space_before = Pt(0)
        
        # Simple list handling
        if line.startswith('- '):
            p.level = 0
            clean_text = line[2:]
            # Main bullet points bold
            run = p.add_run()
            run.text = "• " + clean_text
            set_font(run, Pt(18), True)
            run.font.color.rgb = RGBColor(0, 0, 0)
        elif line.startswith('* '):
            p.level = 1
            clean_text = line[2:]
            run = p.add_run()
            run.text = "  - " + clean_text
            set_font(run, Pt(16), False)
            run.font.color.rgb = RGBColor(80, 80, 80)
        else:
            # Fallback for plain lines
            run = p.add_run()
            run.text = line
            set_font(run, Pt(18))

def generate_ppt():
    prs = Presentation()
    
    # --- Part 2: Work Summary (Projects) ---
    
    # 1. vCube
    add_content_slide(prs, "vCube 平台", """
- 定位：好用的日志在线分析工具 & 统一协作底座
* 核心价值：支持灵活的数据结构定义（文本、数字、人员等），准确映射业务实体；支持多人在线实时编辑，替代传统Excel文件协作，确保数据源唯一可信。
* 关键功能：提供筛选、排序等数据视图，满足定制交付工程师日常维护需求。
* 量化收益：定制交付场景下自动化汇总日报，消除重复劳动；年度总节省工时约 140.6 人天。
""")

    # 2. VLA
    add_content_slide(prs, "VLA (日志分析)", """
- 定位：便捷自由的多维表格工具 & 智能诊断中心
* 核心能力：实现视频帧级播放体验，精准定位卡顿帧；全面解析系统Trace文件，可视化展示耗时瓶颈；支持自定义分析规则与自动化闭环。
* 平台规模：月均处理日志数据 434 TB，成为公司级核心诊断平台。
* 量化收益：通过在线解压与浏览，全年为团队节省约 4,740 人天。
""")

    # 3. Version & FTP
    add_content_slide(prs, "软件版本与 FTP", """
- 定位：便捷使用的系统版本工具 & 高效分发体系
* 核心功能：v8.3.0 重构实现版本上传、自动挂载、安全鉴权、全球分发全流程；云刷机同步功能让工程师本地刷机自动拉取云端镜像。
* 体验优化：实现“一键下载常用刷机包”及 ReleaseNote 在线预览。
* 量化收益：云刷机月使用次数达 12,267 次，年度节省约 9,200 人天。
""")

    # 4. Carrier & Enterprise
    add_content_slide(prs, "运营商与政企定制", """
- 定位：高效高质的定制项目管理工具
* 核心策略：针对运营商与政企项目设计差异化送测与认证流程；开发批量启动送测、一键生成报告及“开在线”自动化功能。
* 质量风控：建立事前拦截机制，自动校验 FCC ID、GMS 认证及安全补丁合规性。
* 量化收益：每年节省约 195 人天；成功拦截 120+ 起认证及合规风险事件。
""")

    # 5. Translation
    add_content_slide(prs, "翻译管理", """
- 定位：高效高质的翻译工具
* 核心创新：AI 赋能上线“智能匹图”功能，利用 AI 自动匹配截图上下文；支持按词条单独驳回，提升协作清晰度。
* 平台扩展：成功扩展至手表产品线，支持全球化需求。
* 量化收益：智能匹图匹配成功率 ~70%；以手表线为例，年化节省达 3,550 人天。
""")

    # 6. Desktop Layout
    add_content_slide(prs, "桌面布局", """
- 概述：微小改动，支撑重大商业成功
* 业务支持：支持对预装文件夹进行批量重命名操作，满足定制需求。
* 价值优化：将 "Recommended Apps" 更名为 "More Apps"，消除多语言显示歧义。
* 量化收益：内部效率节省约 247 人天；外部商业价值方面，成功支撑 INS 预装千万级业务收益。
""")

    # 7. Handover Center
    add_content_slide(prs, "交接中心", """
- 定位：意见交接 MEAT 业务内容 & 协作规范化
* 核心价值：将运营商项目特有的认证信息、定制配置项纳入标准化线上交接流程。
* 风险控制：确保项目转手时关键信息无损传递，降低协作风险。
* 改进方向：持续减少因信息遗漏导致的项目延期事件。
""")

    # 8. Quick App
    add_content_slide(prs, "快应用国家政务服务平台", """
- 年度概况：持续承接国办、部委服务，打造快应用政务一体化平台
* 重点项目：高考专区接入 11 项服务，覆盖 31 个省市，保障高并发稳定性；上线购车、家电等政府补贴查询功能。
* 证照&社保：优化人脸登录体验与安全性，目前证照&社保服务 UV 达 180w+。
* 厂商合作：与华为、OPPO 协作共建快应用政务生态，优化“高效办成一件事”专区。
""")

    # --- Part 3: Problems & Suggestions ---
    add_content_slide(prs, "问题和建议", """
- 遇到的问题
* 接口误报：接口自动化测试因 Nginx 证书过期导致生产环境准确率下降。
* 授权异常：服务器迁移导致系统时间不一致，引发 Token 生成失效。

- 改进建议
* 基础设施监控：建立对证书有效期、系统时间同步的自动化监控与告警机制。
* 提升运维自动化：减少人工配置环节，降低环境迁移带来的风险。
* 持续深化 AI 应用：将已验证成功的 AI 助手推广至更多运维与测试场景。
""")

    # --- Part 4: 2026 Plan ---
    add_content_slide(prs, "2026 工作规划", """
- 智能化深化
* 推广“综合 AI 助手”与“缺陷分析助手”，覆盖更多业务线与研发环节。
* 探索 AI 在自动化代码生成、代码 CR 及智能测试中的深度应用。

- 平台稳定性与性能
* 持续优化政务服务平台，确保高并发场景下的极致体验。
* 提升 vCube 与 VLA 平台的大数据处理能力，应对日志量增长。

- 生态与合作
* 深化与华为、OPPO 等厂商的生态合作，拓展快应用服务边界。

- 团队效能
* 推进工具链整合，打造一站式研发效能平台。
""")

    output_file = "2025_Year_End_Summary.pptx"
    prs.save(output_file)
    print(f"Generate Success: {output_file}")

if __name__ == "__main__":
    generate_ppt()
