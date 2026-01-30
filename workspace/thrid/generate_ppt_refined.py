from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_content_slide(prs, title, work_content, work_insights, key_metrics=None):
    """
    Adds a slide with specific sections: Work Content and Work Insights.
    Optionally adds Key Metrics if provided.
    """
    # Use Title and Content layout (index 1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Body Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    # Helper to add section header
    def add_header(text, color_rgb=(0, 51, 102)): # Dark Blueish
        p = tf.add_paragraph()
        p.text = text
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*color_rgb)
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    # Helper to add bullet point
    def add_bullet(text, level=0):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.level = level
        p.space_after = Pt(4)

    # 1. 工作内容 (Work Content)
    add_header("■ 工作内容")
    for item in work_content:
        add_bullet(item)

    # 2. 关键成效 (Key Metrics - Optional but good for 'summary')
    if key_metrics:
        add_header("■ 关键成效")
        for item in key_metrics:
            add_bullet(item)

    # 3. 工作心得 (Work Insights)
    add_header("■ 工作心得")
    for item in work_insights:
        add_bullet(item)

def main():
    prs = Presentation()
    
    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "2025年度个人工作总结"
    slide.placeholders[1].text = "汇报人：[Your Name]\n部门：[Your Department]"

    # --- Slide 2: Directory ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    tf = slide.placeholders[1].text_frame
    tf.text = "01. 年度工作总结\n02. 问题与建议\n03. 新年工作规划"

    # --- Slide 3: MEAT - vCube ---
    add_content_slide(
        prs, 
        "MEAT平台 - vCube在线协作系统",
        [
            "版本迭代：完成v8.2.9.2至v8.4.4共6个版本，50+次功能迭代。",
            "功能建设：实现文本/选区/成员多维度批注，支持Bug自动流转至Jira。",
            "实时协同：构建多人实时编辑引擎，保证数据唯一性（Single Source of Truth）。"
        ],
        [
            "告别碎片化：从Excel离线交互转向在线实时协作，彻底解决了版本混乱问题。",
            "流程闭环：打通批注到Bug管理的数据链路，让'发现问题'到'解决问题'无缝衔接。",
            "极简主义：工具的价值在于'用完即走'，自动化流转极大降低了用户的操作成本。"
        ],
        [
            "研发/测试节省工时：140.6人天",
            "协作效率提升：90% (消除版本同步耗时)"
        ]
    )

    # --- Slide 4: MEAT - VLA ---
    add_content_slide(
        prs,
        "MEAT平台 - VLA日志分析工具",
        [
            "海量处理：优化大文件加载引擎，支持秒级打开GB级日志。",
            "可视化分析：实现Trace文件时序图展示，标准Log格式化解析。",
            "数据压缩：实施全链路日志压缩策略，降低存储与传输成本。"
        ],
        [
            "数据可视化价值：将枯燥的文本日志转化为直观的时序图，让定位问题从'阅读理解'变为'看图说话'。",
            "性能即体验：在海量数据场景下，加载速度是核心竞争力，压缩算法的优化带来巨大的存储成本节省。"
        ],
        [
            "日志压缩量：434 TB",
            "存储空间折算节省：4,740人天"
        ]
    )

    # --- Slide 5: MEAT - Infrastructure ---
    add_content_slide(
        prs,
        "MEAT平台 - 基础效能建设",
        [
            "版本/FTP优化：实现版本全自动化下载、鉴权与更新检测。",
            "运营商定制：开发自动化校验工具，覆盖FCC ID、GMS认证状态检查。",
            "安全合规：统一桌面布局配置，实施敏感信息拦截策略。"
        ],
        [
            "抓住高频痛点：针对'下载等待'这一高频低价值场景进行优化，产生的ROI（投入产出比）是惊人的。",
            "源头治理：通过自动化工具在代码提交前（Pre-commit）拦截合规问题，远比后期返工高效。"
        ],
        [
            "下载等待时间节省：9,200+人天",
            "合规校验节省：195人天 (拦截120+问题)"
        ]
    )

    # --- Slide 6: MEAT - AI & Handover Center (Replaces old AI slide) ---
    add_content_slide(
        prs,
        "MEAT平台 - AI重构与商业化支撑",
        [
            "AI翻译管理：基于图像识别(CV)技术重构翻译流程，实现UI截图与翻译条目自动匹配。",
            "交接中心：支撑'Recommended Apps'与'More Apps'等INS预装业务的复杂分发逻辑。",
            "商业化闭环：通过微小的配置功能改动，有力支撑了千万级营收的商业项目落地。"
        ],
        [
            "技术赋能业务：AI不仅仅是提效工具，更是解决'人工无法覆盖'场景（如海量截图匹配）的关键。",
            "小功能大价值：交接中心虽然功能点微小，但直接关联公司核心营收，体现了技术对商业成功的直接支撑。"
        ],
        [
            "AI截图匹配节省：1,183人天 (匹配成功率>80%)",
            "交接中心内效提升：节省247人天",
            "商业价值：支撑INS预装千万级营收业务"
        ]
    )

    # --- Slide 7: Quick App Gov Service - Core ---
    add_content_slide(
        prs,
        "快应用 - 国家政务服务平台 (核心保障)",
        [
            "高考/高频保障：圆满完成11省高考查分、31省录取查询保障任务，实现服务零故障。",
            "社保/医保支付：打通医保移动支付全流程，支持多省份社保查询业务。",
            "稳定性建设：建立接口自动化监控体系，解决OAuth鉴权偶发失败问题。"
        ],
        [
            "底线思维：政务服务具有高敏感性，'零故障'是我们坚守的底线。",
            "主动防御：通过自动化监控替代被动客诉，将潜在风险消灭在萌芽状态。"
        ],
        [
            "服务稳定性：99.99%",
            "核心业务保障：覆盖31省高考录取查询"
        ]
    )

    # --- Slide 8: Quick App - Ecosystem & Optimization ---
    add_content_slide(
        prs,
        "快应用 - 厂商合作与体验优化",
        [
            "OPPO生态合作：深度适配OPPO侧卡片样式，部署通用查询接口，提升跨厂商服务体验。",
            "证照/电保专区：重构证照中心展示逻辑，优化交互流程（去除非必要鉴权步骤）。",
            "社保卡专区：优化电子社保卡申领与使用流程，提升卡面加载速度。"
        ],
        [
            "生态共赢：通过技术手段抹平不同厂商间的差异，实现'一次开发，多端运行'的初衷。",
            "体验为王：在证照中心优化中，每一个多余点击的去除，都意味着用户体验的显著提升。"
        ],
        [
            "证照/电保专区流量：PV 90w+, UV 180w+",
            "OPPO侧：成功接入卡片服务与通用查询"
        ]
    )

    # --- Slide 9: Problems & Suggestions ---
    add_content_slide(
        prs,
        "问题回顾与改进建议",
        [
            "技术债：部分老旧系统接口缺乏监控，OAuth Code依赖系统时间同步导致偶发失败。",
            "流程瓶颈：人工Code Review在高峰期效率低，且容易遗漏基础规范问题。"
        ],
        [
            "系统性治理：不能头痛医头，建议推动运维侧统一NTP时间同步，并建立全链路Nginx配置标准。",
            "工具化思维：建议推广AI辅助Code Review工具，将规范检查自动化，让人聚焦在逻辑架构审查上。"
        ]
    )

    # --- Slide 9: New Year Plan ---
    add_content_slide(
        prs,
        "2026年工作规划",
        [
            "业务深耕：落地医保移动支付全流程，打造政务服务标杆案例。",
            "平台进化：MEAT平台引入更多AI Agent（如自动生成测试用例），进一步降低研发成本。",
            "技术沉淀：输出高质量前端组件库，探索鸿蒙原生应用开发，保持技术栈先进性。"
        ],
        [
            "从'支撑'到'驱动'：技术团队不能仅满足于支撑业务，更要通过平台化、智能化工具反向驱动业务效率提升。",
            "持续学习：拥抱AI 2.0时代，让每个工程师都成为'超级个体'。"
        ]
    )

    # --- Slide 10: End ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "谢谢聆听"
    slide.placeholders[1].text = "请批评指正"

    output_path = os.path.join(os.getcwd(), "2025_Year_End_Summary_Refined_v2.pptx")
    prs.save(output_path)
    print(f"Refined PPTX generated at: {output_path}")

if __name__ == "__main__":
    main()
