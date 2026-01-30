from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_slide(prs, title, work_content, work_insights, key_metrics=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    slide.shapes.title.text = title
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    def add_header(text, color_rgb=(0, 51, 102)):
        p = tf.add_paragraph()
        p.text = text
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*color_rgb)
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    def add_bullet(text, level=0):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.level = level
        p.space_after = Pt(4)

    # 1. 工作内容
    add_header("■ 工作内容")
    for item in work_content:
        add_bullet(item)

    # 2. 关键成效
    if key_metrics:
        add_header("■ 关键成效")
        for item in key_metrics:
            add_bullet(item)

    # 3. 工作心得
    add_header("■ 工作心得")
    for item in work_insights:
        add_bullet(item)

def main():
    prs = Presentation()
    
    # Slide 1: Directory (As Title/End are skipped/manual)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    tf = slide.placeholders[1].text_frame
    tf.text = "01. 年度工作总结 (8个项目)\n02. 问题与建议\n03. 新年工作规划"

    # --- Project 1: vCube ---
    create_slide(prs, "vCube", 
        [
            "版本迭代：完成v8.2.9.2至v8.4.4共6个版本，50+次功能迭代。",
            "统一协作：实现文本/选区/成员多类型批注，Bug自动流转。",
            "实时协同：多人实时编辑，保证数据唯一性，替代Excel离线流转。"
        ],
        [
            "告别Excel版本混乱，实现数据源头的统一与实时同步。",
            "工具价值在于'用完即走'，自动化闭环降低了协作成本。"
        ],
        ["研发/测试节省：140.6人天", "协作效率提升：90%"]
    )

    # --- Project 2: VLA ---
    create_slide(prs, "VLA", 
        [
            "海量日志：优化大文件加载引擎，支持GB级日志秒级打开。",
            "可视化：实现Trace文件时序图展示，标准Log格式化解析。",
            "全链路压缩：实施日志压缩策略，大幅降低存储传输成本。"
        ],
        [
            "将枯燥文本转化为可视化图表，极大降低了问题定位门槛。",
            "性能优化直接转化为体验提升，秒级打开大文件成为核心竞争力。"
        ],
        ["日志压缩量：434 TB", "存储节省折算：4,740人天"]
    )

    # --- Project 3: 运营商与政企定制 ---
    create_slide(prs, "运营商与政企定制", 
        [
            "自动化校验：开发工具自动校验FCC ID、GMS认证状态。",
            "源头治理：在代码提交前拦截合规问题，避免后期返工。",
            "双端适配：解决运营商项目不同配置的兼容性问题。"
        ],
        [
            "用工具替代人工CheckList，不仅快，而且绝对准确。",
            "合规是红线，自动化是守住红线的最高效手段。"
        ],
        ["拦截合规问题：120+项", "人工校验节省：195人天"]
    )

    # --- Project 4: 软件版本与FTP ---
    create_slide(prs, "软件版本与FTP", 
        [
            "自动化下载：实现版本全自动下载、鉴权与更新检测。",
            "ReleaseNote解析：自动提取版本信息，减少人工整理时间。",
            "高频场景优化：针对每日高频使用的下载场景进行极致提速。"
        ],
        [
            "抓住'高频低价值'场景（如等待下载）进行优化，ROI极高。",
            "极简工具链让工程师能更专注于核心业务开发。"
        ],
        ["下载等待节省：9,200+人天", "使用频率：12,000+次"]
    )

    # --- Project 5: 翻译管理 ---
    create_slide(prs, "翻译管理", 
        [
            "AI截图匹配：基于CV技术实现UI截图与翻译条目自动匹配。",
            "流程重构：替代传统人工截图与匹配流程，大幅提升效率。",
            "质量把控：自动检测翻译缺失与错误，提升国际化质量。"
        ],
        [
            "AI在'重复性劳动'场景下的最佳实践，释放人力。",
            "匹配成功率>80%是工具从'可用'到'好用'的分水岭。"
        ],
        ["截图匹配节省：1,183人天", "匹配成功率：>80%"]
    )

    # --- Project 6: 桌面布局 ---
    create_slide(prs, "桌面布局", 
        [
            "统一标准：规范桌面布局配置，确保企业视觉统一性。",
            "安全合规：实施敏感信息拦截策略，防止数据泄露。",
            "风险防控：通过技术手段规避法律与舆情风险。"
        ],
        [
            "看似微小的改动，实则支撑了重大的企业安全合规需求。",
            "安全无小事，防患于未然是最高级的安全策略。"
        ],
        ["安全风险拦截：100%覆盖", "合规事故：0起"]
    )

    # --- Project 7: 交接中心 ---
    create_slide(prs, "交接中心", 
        [
            "业务支撑：支持INS预装、Recommended Apps等复杂分发业务。",
            "配置灵活：通过微小改动支撑千万级营收项目的落地。",
            "流程优化：简化交接流程，提升跨部门协作效率。"
        ],
        [
            "技术直接支撑商业成功，小功能撬动大收益。",
            "打通业务流转的'最后一公里'，确保商业变现畅通无阻。"
        ],
        ["内部协作节省：247人天", "支撑营收：千万级"]
    )

    # --- Project 8: 快应用国家政务服务平台 ---
    create_slide(prs, "快应用国家政务服务平台", 
        [
            "核心保障：零故障支撑高考查分、社保医保支付等高频服务。",
            "生态合作：适配OPPO侧卡片，部署通用查询接口。",
            "体验重构：优化证照中心与电子社保卡流程，去除非必要鉴权。"
        ],
        [
            "政务服务'零故障'是底线，必须保持对生产环境的敬畏。",
            "每一个多余点击的去除，都是对用户体验的极致追求。"
        ],
        ["服务稳定性：99.99%", "核心流量：PV 90w+, UV 180w+"]
    )

    # --- Part 3: Problems & Suggestions ---
    create_slide(prs, "问题与建议", 
        [
            "技术痛点：OAuth Code依赖系统时间，Nginx鉴权偶发失败。",
            "流程瓶颈：人工Code Review效率低，基础规范依赖人工检查。"
        ],
        [
            "系统治理：推进全链路NTP时间同步与Nginx配置标准化。",
            "工具赋能：推广AI辅助Code Review，实现规范检查自动化。"
        ],
        []
    )

    # --- Part 4: New Year Plan ---
    create_slide(prs, "新年工作规划", 
        [
            "快应用业务：落地医保移动支付，打造政务服务标杆。",
            "平台深化：MEAT平台引入更多AI Agent，降低研发成本。",
            "技术沉淀：输出高质量组件库，探索鸿蒙原生应用开发。"
        ],
        [
            "从支撑向驱动转型，用技术反哺业务增长。",
            "拥抱AI 2.0，构建人机协同的新型研发模式。"
        ],
        []
    )

    output_path = os.path.join(os.getcwd(), "2025_Year_End_Summary_Final.pptx")
    prs.save(output_path)
    print(f"Final PPTX generated at: {output_path}")

if __name__ == "__main__":
    main()
