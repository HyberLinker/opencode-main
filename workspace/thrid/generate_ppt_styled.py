from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

# --- Configuration ---
THEME_COLOR_BLUE = RGBColor(0, 80, 158)   # Corporate Blue
THEME_COLOR_RED = RGBColor(220, 50, 50)   # Highlight Red
TEXT_COLOR_MAIN = RGBColor(50, 50, 50)    # Dark Grey
TEXT_COLOR_LIGHT = RGBColor(100, 100, 100) # Lighter Grey
BG_COLOR_LIGHT = RGBColor(245, 247, 250)  # Very light blue/grey for boxes

FONT_CN = "HarmonyOS Sans SC"
FONT_EN = "Inter"

def set_font(paragraph, size=Pt(14), bold=False, color=TEXT_COLOR_MAIN, is_title=False):
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = FONT_EN # Set English/Number font
    
    # Set Chinese font using OXML workaround
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    rPr = run._r.get_or_add_rPr()
    ea = qn('a:ea')
    rPr.set(ea, FONT_CN)

def create_problems_slide(prs, title, problems_list):
    """
    Custom layout for Problems & Suggestions:
    Uses horizontal cards to map Problem -> Impact -> Solution
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    # Title
    create_title_bar(slide, title)
    
    # Grid Settings
    start_top = Inches(1.5)
    card_height = Inches(2.2)
    gap = Inches(0.3)
    
    # Columns: Problem (30%), Impact (30%), Solution (40%)
    col1_left = Inches(0.5)
    col1_width = Inches(3.5)
    
    col2_left = col1_left + col1_width + Inches(0.2)
    col2_width = Inches(3.5)
    
    col3_left = col2_left + col2_width + Inches(0.2)
    col3_width = Inches(4.5)

    # Headers
    headers = [("痛点与挑战", col1_left), ("负面影响", col2_left), ("改进方案", col3_left)]
    for text, left in headers:
        txBox = slide.shapes.add_textbox(left, Inches(1.1), Inches(3), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        set_font(p, size=Pt(16), bold=True, color=THEME_COLOR_BLUE)

    current_top = start_top
    
    for prob in problems_list:
        # Draw background container for the row
        bg_shape = slide.shapes.add_shape(
            1, # Rectangle
            Inches(0.4), current_top, Inches(12.5), card_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR_LIGHT
        bg_shape.line.fill.background()

        # 1. Problem
        txBox = slide.shapes.add_textbox(col1_left, current_top + Inches(0.2), col1_width, card_height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = prob['title']
        set_font(p, size=Pt(14), bold=True, color=THEME_COLOR_RED)
        
        p = tf.add_paragraph()
        p.text = prob['desc']
        set_font(p, size=Pt(12), color=TEXT_COLOR_MAIN)

        # 2. Impact
        txBox = slide.shapes.add_textbox(col2_left, current_top + Inches(0.2), col2_width, card_height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = prob['impact']
        set_font(p, size=Pt(12), color=TEXT_COLOR_MAIN)

        # 3. Solution
        txBox = slide.shapes.add_textbox(col3_left, current_top + Inches(0.2), col3_width, card_height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = prob['solution']
        set_font(p, size=Pt(13), bold=True, color=THEME_COLOR_BLUE)
        
        current_top += card_height + gap

def create_plan_slide(prs, title, plan_modules):
    """
    Custom layout for New Year Plan:
    Uses 3 vertical columns for different dimensions.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    # Title
    create_title_bar(slide, title)
    
    # Layout Settings
    start_left = Inches(0.5)
    gap = Inches(0.4)
    col_width = Inches(3.8)
    top = Inches(1.5)
    height = Inches(5.0)
    
    for i, module in enumerate(plan_modules):
        current_left = start_left + (i * (col_width + gap))
        
        # 1. Header Box
        header_shape = slide.shapes.add_shape(
            1, # Rectangle
            current_left, top, col_width, Inches(0.6)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = THEME_COLOR_BLUE
        header_shape.line.fill.background()
        
        tf = header_shape.text_frame
        p = tf.paragraphs[0]
        p.text = module['title']
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=Pt(18), bold=True, color=RGBColor(255, 255, 255))
        
        # 2. Content Box
        content_shape = slide.shapes.add_shape(
            1, # Rectangle
            current_left, top + Inches(0.6), col_width, height
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = BG_COLOR_LIGHT
        content_shape.line.fill.background()
        
        tf = content_shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.2)
        
        for item in module['items']:
            # Sub-title
            p = tf.add_paragraph()
            p.text = "• " + item['sub']
            set_font(p, size=Pt(14), bold=True, color=TEXT_COLOR_MAIN)
            p.space_before = Pt(12)
            
            # Detail
            p = tf.add_paragraph()
            p.text = item['detail']
            set_font(p, size=Pt(12), color=TEXT_COLOR_LIGHT)
            p.level = 1

def create_title_bar(slide, title):
    """Helper to create the standard title bar"""
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(10)
    height = Inches(0.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=Pt(28), bold=True, color=THEME_COLOR_BLUE)
    
    # Decorative Line
    shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR_BLUE
    shape.line.fill.background()

def create_styled_slide(prs, title, work_content, work_insights, key_metrics=None):
    # Use Blank Layout to manually position everything
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    # Title
    create_title_bar(slide, title)

    # ... (Rest of the function remains similar, just adjusted to use create_title_bar helper implicitly or logic copied)
    # --- 2. Left Column: Work Content (60% width) ---
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(7.5)
    height = Inches(5.0)
    
    # Section Header
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "■ 工作内容"
    set_font(p, size=Pt(18), bold=True, color=THEME_COLOR_BLUE)
    
    # Content Bullets
    top += Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for item in work_content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        set_font(p, size=Pt(16), color=TEXT_COLOR_MAIN)

    # --- 3. Right Column: Metrics & Insights (30% width) ---
    bg_left = Inches(8.5)
    bg_top = Inches(1.8)
    bg_width = Inches(4.3)
    bg_height = Inches(5.0)
    
    # Metrics Section
    if key_metrics:
        txBox = slide.shapes.add_textbox(bg_left, bg_top, bg_width, Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = "■ 关键成效"
        set_font(p, size=Pt(18), bold=True, color=THEME_COLOR_RED)
        
        curr_top = bg_top + Inches(0.5)
        for metric in key_metrics:
            if "：" in metric:
                label, value = metric.split("：", 1)
                txBox = slide.shapes.add_textbox(bg_left, curr_top, bg_width, Inches(0.3))
                p = txBox.text_frame.paragraphs[0]
                p.text = label
                set_font(p, size=Pt(12), color=TEXT_COLOR_LIGHT)
                
                curr_top += Inches(0.3)
                txBox = slide.shapes.add_textbox(bg_left, curr_top, bg_width, Inches(0.5))
                p = txBox.text_frame.paragraphs[0]
                p.text = value
                set_font(p, size=Pt(22), bold=True, color=THEME_COLOR_RED)
                curr_top += Inches(0.5)
            else:
                txBox = slide.shapes.add_textbox(bg_left, curr_top, bg_width, Inches(0.4))
                p = txBox.text_frame.paragraphs[0]
                p.text = metric
                set_font(p, size=Pt(14), bold=True, color=THEME_COLOR_RED)
                curr_top += Inches(0.4)
        bg_top = curr_top + Inches(0.3)

    # Insights Section
    txBox = slide.shapes.add_textbox(bg_left, bg_top, bg_width, Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "■ 工作心得"
    set_font(p, size=Pt(18), bold=True, color=THEME_COLOR_BLUE)
    
    curr_top = bg_top + Inches(0.5)
    txBox = slide.shapes.add_textbox(bg_left, curr_top, bg_width, Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for item in work_insights:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_font(p, size=Pt(14), color=TEXT_COLOR_MAIN)

def create_directory_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(5), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "目录"
    set_font(p, size=Pt(28), bold=True, color=THEME_COLOR_BLUE)

    # Content
    content = [
        "01. 年度工作总结",
        "    1.1 vCube 协作平台",
        "    1.2 VLA 日志分析",
        "    1.3 基础效能工具",
        "    1.4 AI 创新应用",
        "    1.5 快应用政务服务",
        "02. 问题回顾与建议",
        "03. 新年工作规划"
    ]
    
    left = Inches(1.5)
    top = Inches(2.0)
    for line in content:
        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = line
        size = Pt(20) if not line.startswith("    ") else Pt(16)
        bold = not line.startswith("    ")
        set_font(p, size=size, bold=bold, color=TEXT_COLOR_MAIN)
        top += Inches(0.6) if not line.startswith("    ") else Inches(0.4)

def main():
    prs = Presentation()
    # Set slide size to 16:9 just in case (default is usually 4:3 in some versions, but 16:9 in newer)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- Directory ---
    create_directory_slide(prs)

    # --- Project 1: vCube ---
    create_styled_slide(prs, "vCube 协作平台", 
        [
            "版本迭代：完成 v8.2.9.2 - v8.4.4 共6个版本，50+次功能迭代。",
            "统一协作：实现文本、选区、成员等多种批注类型，支持Bug一键流转Jira。",
            "实时协同：构建多人实时编辑引擎，保证数据一致性，替代传统Excel离线交互模式。"
        ],
        [
            "价值主张：工具的价值在于'用完即走'，自动化流转降低了用户的操作成本。",
            "协作变革：彻底消除Excel版本混乱，实现数据源头的统一。"
        ],
        ["研发测试节省：140.6人天", "协作效率提升：90%"]
    )

    # --- Project 2: VLA ---
    create_styled_slide(prs, "VLA 日志分析工具", 
        [
            "海量日志引擎：优化大文件加载算法，支持GB级日志秒级打开。",
            "深度可视化：实现Trace文件时序图展示，支持标准Log格式化解析与过滤。",
            "全链路压缩：实施日志上传与存储压缩策略，显著降低成本。"
        ],
        [
            "体验升级：将枯燥的文本阅读转化为直观的'看图说话'，降低问题定位门槛。",
            "性能即体验：秒级打开大文件是核心竞争力。"
        ],
        ["日志压缩量：434 TB", "存储节省折算：4,740人天"]
    )

    # --- Project 3: Customization ---
    create_styled_slide(prs, "运营商与政企定制", 
        [
            "自动化校验：开发工具自动检查FCC ID、GMS认证状态，覆盖双端配置。",
            "源头治理：在代码提交阶段拦截合规问题，避免版本回退。",
            "兼容适配：解决运营商项目复杂的配置差异问题。"
        ],
        [
            "合规红线：用工具替代人工CheckList，确保100%准确。",
            "防患未然：在源头解决问题成本最低。"
        ],
        ["拦截合规问题：120+项", "人工校验节省：195人天"]
    )

    # --- Project 4: Version & FTP ---
    create_styled_slide(prs, "软件版本与FTP工具", 
        [
            "全流程自动化：实现版本自动下载、鉴权校验与更新检测。",
            "信息提取：自动解析ReleaseNote，减少人工整理工作量。",
            "场景优化：针对每日高频使用的下载场景进行极致提速。"
        ],
        [
            "ROI思维：抓住'高频低价值'场景（如等待下载）进行优化，收益巨大。",
            "极简工具链：让工程师专注于核心创造性工作。"
        ],
        ["下载等待节省：9,200+人天", "使用频率：12,000+次"]
    )

    # --- Project 5: Translation & AI ---
    create_styled_slide(prs, "翻译管理 (AI重构)", 
        [
            "AI截图匹配：基于CV技术实现UI截图与翻译条目自动关联，匹配率>80%。",
            "流程自动化：替代人工截图与查找流程，大幅提升国际化适配效率。",
            "质量监控：自动检测翻译缺失与错误，保障多语言版本质量。"
        ],
        [
            "AI务实落地：聚焦'重复性劳动'场景，释放人力资源。",
            "信任建立：只有高准确率才能让用户真正依赖AI工具。"
        ],
        ["截图匹配节省：1,183人天", "匹配成功率：>80%"]
    )

    # --- Project 6: Desktop Layout ---
    create_styled_slide(prs, "桌面布局规范", 
        [
            "统一标准：制定并固化桌面布局配置，确保企业视觉风格统一。",
            "安全合规：实施敏感信息拦截策略，防止数据意外泄露。",
            "风险管控：通过技术手段规避潜在的法律与舆情风险。"
        ],
        [
            "细节决定成败：微小的规范化改动支撑了企业的宏观安全战略。",
            "零事故目标：安全工作没有侥幸，必须100%覆盖。"
        ],
        ["安全风险拦截：100%覆盖", "合规事故：0起"]
    )

    # --- Project 7: Handover Center ---
    create_styled_slide(prs, "交接中心", 
        [
            "业务支撑：支持INS预装、Recommended Apps等复杂分发业务逻辑。",
            "配置灵活：通过低成本配置改动，支撑千万级营收项目快速落地。",
            "流程优化：简化跨部门交接流程，提升业务流转效率。"
        ],
        [
            "商业价值：技术直接服务于营收，小功能撬动大收益。",
            "敏捷响应：快速响应商业需求是技术团队的核心价值之一。"
        ],
        ["内部协作节省：247人天", "支撑营收：千万级"]
    )

    # --- Project 8: Quick App Gov ---
    create_styled_slide(prs, "快应用国家政务服务平台", 
        [
            "核心保障：零故障支撑高考查分、社保医保支付等亿级流量服务。",
            "生态连接：适配OPPO侧卡片样式，部署通用查询接口，打破厂商壁垒。",
            "体验重构：优化证照中心交互，去除非必要鉴权，提升转化率。"
        ],
        [
            "底线思维：政务服务关乎民生，稳定性是最高优先级。",
            "体验为王：每一个多余点击的去除，都是对用户的尊重。"
        ],
        ["服务稳定性：99.99%", "核心流量：PV 90w+"]
    )

    # --- Part 3: Problems & Suggestions ---
    create_problems_slide(prs, "问题回顾与改进建议", 
        [
            {
                "title": "基础设施稳定性",
                "desc": "OAuth Code依赖系统时间同步，Nginx鉴权在系统时间漂移时偶发失败。",
                "impact": "导致用户登录失败，引发客诉；部分旧系统接口缺乏秒级监控，响应滞后。",
                "solution": "1. 推进运维侧统一NTP时间同步，消除时间漂移隐患。\n2. 建立全链路Nginx配置标准与秒级监控报警。"
            },
            {
                "title": "研发效能瓶颈",
                "desc": "Code Review主要依赖人工，在发版高峰期效率低下。",
                "impact": "基础规范问题容易遗漏，不仅占用高级人力，还可能导致线上隐患。",
                "solution": "引入AI辅助Code Review工具，将命名规范、代码风格等检查自动化，让人聚焦逻辑架构。"
            }
        ]
    )

    # --- Part 4: New Year Plan ---
    create_plan_slide(prs, "2026年工作规划", 
        [
            {
                "title": "业务深耕",
                "items": [
                    {"sub": "医保移动支付", "detail": "打通医保支付全流程，落地标杆省份，实现政务服务闭环。"},
                    {"sub": "政务服务标杆", "detail": "优化证照中心体验，提升用户转化率，打造行业样板。"}
                ]
            },
            {
                "title": "平台进化",
                "items": [
                    {"sub": "MEAT Agent", "detail": "引入智能体自动生成测试用例，降低QA回归成本。"},
                    {"sub": "日志分析升级", "detail": "VLA支持更多日志格式解析，覆盖更多业务场景。"}
                ]
            },
            {
                "title": "技术沉淀",
                "items": [
                    {"sub": "组件库建设", "detail": "沉淀高质量通用前端组件库，提升开发效率与UI一致性。"},
                    {"sub": "鸿蒙原生探索", "detail": "预研鸿蒙Next开发技术，储备原生应用开发能力。"}
                ]
            }
        ]
    )

    output_path = os.path.join(os.getcwd(), "2025_Year_End_Summary_Final_Styled.pptx")
    prs.save(output_path)
    print(f"Styled PPTX generated at: {output_path}")

if __name__ == "__main__":
    main()
