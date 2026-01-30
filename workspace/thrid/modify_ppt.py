from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

# --- Style Config (Consistent with previous script) ---
THEME_COLOR_BLUE = RGBColor(0, 80, 158)
TEXT_COLOR_MAIN = RGBColor(50, 50, 50)
FONT_CN = "HarmonyOS Sans SC"
FONT_EN = "Inter"

def set_font(paragraph, size=Pt(14), bold=False, color=TEXT_COLOR_MAIN):
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = FONT_EN 
    
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    rPr = run._r.get_or_add_rPr()
    ea = qn('a:ea')
    rPr.set(ea, FONT_CN)

def update_handover_slide(filepath):
    prs = Presentation(filepath)
    
    target_slide = None
    for i, slide in enumerate(prs.slides):
        # iterate shapes to find title
        for shape in slide.shapes:
            if shape.has_text_frame and "交接中心" in shape.text:
                # Double check it's likely a title (top of page)
                if shape.top < Pt(100): 
                    target_slide = slide
                    print(f"Found '交接中心' slide at index {i}")
                    break
        if target_slide:
            break
            
    if not target_slide:
        print("未找到'交接中心'页面")
        return

    print("Found '交接中心' slide. Updating content...")

    # New Content
    new_work_content = [
        "多业务支撑：支持 INS预装、Recommended Apps、More Apps 等多种分发业务的差异化交接。",
        "精细化管控：实现 敏感应用屏蔽 及 商业化/测试标签 的动态配置，确保分发合规。",
        "商业化闭环：通过微小的配置能力升级，打通了从预装到变现的 最后一公里。"
    ]
    
    new_insights = [
        "小支点大杠杆：交接配置虽然只是微小的功能点，但它是千万级营收项目落地的 关键依赖。",
        "安全即收益：通过精准的合规拦截（如屏蔽敏感应用），在保障营收的同时规避了潜在的 运营风险。"
    ]

    # Iterate shapes to find the content boxes
    # Heuristic: Find text box containing keywords from OLD content
    
    content_updated = False
    insights_updated = False

    for shape in target_slide.shapes:
        if not shape.has_text_frame:
            continue
            
        text = shape.text
        
        # 1. Update Work Content
        if "业务支撑" in text and "配置灵活" in text:
            tf = shape.text_frame
            tf.clear()
            for item in new_work_content:
                p = tf.add_paragraph()
                p.text = item
                p.space_after = Pt(10)
                set_font(p, size=Pt(16), color=TEXT_COLOR_MAIN)
            content_updated = True
            print("Updated Work Content.")

        # 2. Update Insights
        elif "商业价值" in text and "敏捷响应" in text:
            tf = shape.text_frame
            tf.clear()
            for item in new_insights:
                p = tf.add_paragraph()
                p.text = item
                p.space_after = Pt(8)
                set_font(p, size=Pt(14), color=TEXT_COLOR_MAIN)
            insights_updated = True
            print("Updated Insights.")

    if content_updated and insights_updated:
        output_path = filepath # Overwrite or save as new
        prs.save(output_path)
        print(f"Successfully updated: {output_path}")
    else:
        print("Could not locate specific text boxes to update. Please check if content matches expected patterns.")

if __name__ == "__main__":
    ppt_file = os.path.join(os.getcwd(), "2025_Year_End_Summary_Final_Styled.pptx")
    update_handover_slide(ppt_file)
