import os
from pptx import Presentation
from docx import Document

def extract_pptx_info(filepath):
    print(f"--- Analyzing PPTX: {filepath} ---")
    if not os.path.exists(filepath):
        print("File not found.")
        return

    prs = Presentation(filepath)
    print(f"Total Slides: {len(prs.slides)}")
    print(f"Slide Width: {prs.slide_width}, Height: {prs.slide_height}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1} Layout: {slide.slide_layout.name}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  - Text: {shape.text[:50]}..." if len(shape.text) > 50 else f"  - Text: {shape.text}")

def extract_docx_info(filepath):
    print(f"\n--- Extracting DOCX: {filepath} ---")
    if not os.path.exists(filepath):
        print("File not found.")
        return

    doc = Document(filepath)
    for para in doc.paragraphs:
        if para.text.strip():
            print(f"Para: {para.text}")

base_skill_dir = r"C:\Users\11101526\.config\opencode\skills\年终总结ppt-skill"
example_ppt = os.path.join(base_skill_dir, "reference", "output example", "个人年度工作总结--刘家玮.pptx")
input_doc1 = os.path.join(base_skill_dir, "reference", "input", "25年年终总结.docx")
input_doc2 = os.path.join(base_skill_dir, "reference", "input", "快应用政务服务-2025工作总结.docx")

extract_pptx_info(example_ppt)
extract_docx_info(input_doc1)
extract_docx_info(input_doc2)
