from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

# --- Constants ---
TEXT_COLOR_MAIN = RGBColor(50, 50, 50)
TEXT_COLOR_LIGHT = RGBColor(100, 100, 100)
THEME_COLOR_RED = RGBColor(220, 50, 50)
FONT_CN = "HarmonyOS Sans SC"
FONT_EN = "Inter"

def set_font(paragraph, size=Pt(14), bold=False, color=TEXT_COLOR_MAIN):
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = FONT_EN 
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    rPr = run._r.get_or_add_rPr()
    ea = qn('a:ea')
    rPr.set(ea, FONT_CN)

def update_text_frame(shape, text_items, size=Pt(16), is_metric=False):
    tf = shape.text_frame
    tf.clear()
    if is_metric:
        # Single line metric value
        p = tf.paragraphs[0]
        p.text = text_items[0]
        set_font(p, size=size, bold=True, color=THEME_COLOR_RED)
    else:
        # Bullet list
        for item in text_items:
            p = tf.add_paragraph()
            p.text = item
            p.space_after = Pt(10)
            set_font(p, size=size, color=TEXT_COLOR_MAIN)

def process_ppt(filepath):
    print(f"Opening {filepath}...")
    prs = Presentation(filepath)
    
    slide_handover = None
    slide_desktop = None
    
    # 1. Locate Slides
    for slide in prs.slides:
        # Heuristic: Check title or top text
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top < Inches(1.5):
                if "交接中心" in shape.text:
                    slide_handover = slide
                elif "桌面布局" in shape.text:
                    slide_desktop = slide
        
    # 2. Update 交接中心 (Handover)
    if slide_handover:
        print("Found '交接中心' slide. Updating...")
        # New Content
        content_handover = [
            "流程标准化：将运营商项目特有的认证信息、定制配置项纳入标准化线上交接流程。",
            "资产管理：统一管理关键项目资产，确保项目转手时信息的完整性与安全性。",
            "风险防控：消除因信息遗漏导致的项目延期风险。"
        ]
        insights_handover = [
            "无损传递：信息的完整交接是项目平稳过渡的基石，降低了跨团队协作的摩擦成本。",
            "流程即防线：标准化的交接流程是规避人为失误的最有效防线。"
        ]
        
        shapes_to_delete = []
        
        for shape in slide_handover.shapes:
            if not shape.has_text_frame: continue
            
            # Left Column (Content) - Heuristic: Left < 5 inches, Top > 2 inches
            if shape.left < Inches(5) and shape.top > Inches(2):
                update_text_frame(shape, content_handover, size=Pt(16))
                
            # Right Column (Insights) - Heuristic: Left > 8 inches, Top > 3.5 inches (Insights are lower)
            elif shape.left > Inches(8) and shape.top > Inches(3.5):
                 update_text_frame(shape, insights_handover, size=Pt(14))
            
            # Right Column (Metrics to Remove) - Heuristic: Left > 8 inches, Top between 2 and 3.5 inches
            # Also check text content to be sure we don't delete headers
            elif shape.left > Inches(8) and shape.top < Inches(3.5) and shape.top > Inches(1.5):
                text = shape.text
                if "关键成效" in text:
                    pass # Keep header? Or remove if no metrics? 
                    # If no metrics, we should probably remove "关键成效" header too.
                    shapes_to_delete.append(shape)
                elif "247" in text or "千万" in text or "内部" in text or "营收" in text or "INS" in text:
                     shapes_to_delete.append(shape)
                elif "节省" in text or "提升" in text: # Catch labels
                     shapes_to_delete.append(shape)

        # Remove metrics shapes
        for shape in shapes_to_delete:
            sp = shape._element
            sp.getparent().remove(sp)
            print("Removed metric shape from Handover slide.")

    # 3. Update 桌面布局 (Desktop)
    if slide_desktop:
        print("Found '桌面布局' slide. Updating...")
        content_desktop = [
            "批量处理：支持对预装文件夹进行批量重命名，满足区域化定制需求。",
            "歧义消除：将 'Recommended Apps' 更名为中性的 'More Apps'，解决多语言显示歧义。",
            "商业支撑：快速响应 INS 预装业务需求，确保商务合同如期落地。"
        ]
        insights_desktop = [
            "微小改动，重大价值：一个文件夹名称的修改，直接撬动了千万级别的商业合作。",
            "敏捷响应：技术平台对前端业务的直接驱动力，体现在对商业需求的快速落地能力上。"
        ]
        
        for shape in slide_desktop.shapes:
            if not shape.has_text_frame: continue
            text = shape.text
            
            # Left Column (Content)
            if shape.left < Inches(5) and shape.top > Inches(2):
                update_text_frame(shape, content_desktop, size=Pt(16))
            
            # Right Column (Insights)
            elif shape.left > Inches(8) and shape.top > Inches(4): # Insights are usually at bottom
                 update_text_frame(shape, insights_desktop, size=Pt(14))
            
            # Right Column (Metrics) - We reuse the boxes
            # Current: "安全风险拦截" -> "内部节省", "100%覆盖" -> "247人天"
            # Current: "合规事故" -> "支撑营收", "0起" -> "千万级"
            elif shape.left > Inches(8) and shape.top < Inches(4) and shape.top > Inches(1.5):
                if "安全" in text and "风险" in text:
                    shape.text_frame.text = "内部节省"
                elif "100%" in text or "覆盖" in text:
                    update_text_frame(shape, ["247人天"], size=Pt(22), is_metric=True)
                elif "合规" in text or "事故" in text:
                    shape.text_frame.text = "支撑营收"
                elif "0起" in text:
                    update_text_frame(shape, ["千万级"], size=Pt(22), is_metric=True)

    output_file = filepath.replace(".pptx", "_Corrected.pptx")
    prs.save(output_file)
    print(f"Completed! Saved to {output_file}")

if __name__ == "__main__":
    target_file = os.path.join(os.getcwd(), "2025_Year_End_Summary_Final_Styled.pptx")
    process_ppt(target_file)
