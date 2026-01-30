from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_slide(prs, title_text, content_items):
    # Use a blank layout (usually index 6 in default template) or Title and Content (index 1)
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Content
    # content_items is a list of dictionaries: {"header": str, "body": str}
    # We will format them into the body placeholder
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear default text
    
    for i, item in enumerate(content_items):
        p = tf.add_paragraph()
        p.text = item["header"]
        p.font.bold = True
        p.font.size = Pt(20)
        p.level = 0
        
        if item["body"]:
            p_body = tf.add_paragraph()
            p_body.text = item["body"]
            p_body.font.size = Pt(16)
            p_body.level = 1
        
        # Add spacing
        if i < len(content_items) - 1:
            p_spacer = tf.add_paragraph()
            p_spacer.text = ""
            p_spacer.font.size = Pt(5)

def create_title_slide(prs, main_title, sub_title):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = main_title
    slide.placeholders[1].text = sub_title

def main():
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(prs, "2025年个人年终总结", "高效 协同 极简\n汇报人：[Your Name]")
    
    # Slide 2: Agenda
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "目录"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. 工作总结"
    p = tf.add_paragraph()
    p.text = "2. 问题与建议"
    p = tf.add_paragraph()
    p.text = "3. 新年工作规划"

    # Slide 3: MEAT - vCube
    create_slide(prs, "MEAT应用 - vCube平台与统一协作", [
        {"header": "背景介绍", "body": "vCube平台实现统一协作，覆盖6个版本，50+次迭代。解决多工具割裂问题。"},
        {"header": "工作亮点", "body": "支持文本/选区/成员多类型批注；Bug自动流转；实时协同编辑。"},
        {"header": "核心价值", "body": "研发/测试节省人天 140.6天；消除Excel版本混乱，提升协作效率90%。"}
    ])

    # Slide 4: MEAT - VLA
    create_slide(prs, "MEAT应用 - VLA日志分析", [
        {"header": "背景介绍", "body": "解决海量日志查看难、下载慢的问题；支持Trace文件在线可视化。"},
        {"header": "工作亮点", "body": "强力压缩日志（434TB数据处理）；自定义过滤与时序图展示；秒级打开超大文件。"},
        {"header": "核心价值", "body": "节省存储空间折算约4740人天；大幅提升问题定位效率，减少等待时间。"}
    ])

    # Slide 5: MEAT - Base & Customization
    create_slide(prs, "MEAT应用 - 基础建设与定制化", [
        {"header": "版本与FTP优化", "body": "实现全自动化下载/鉴权；节省下载等待时间折算9200+人天（高频使用场景）。"},
        {"header": "运营商定制化", "body": "自动化校验FCC/GMS认证状态；拦截120+合规问题；节省人工校验195人天。"},
        {"header": "桌面布局规范", "body": "统一桌面布局配置，防止敏感信息泄露，提升企业安全合规。"}
    ])

    # Slide 6: MEAT - AI & Innovation
    create_slide(prs, "MEAT应用 - AI与效能创新", [
        {"header": "AI翻译管理", "body": "基于AI的UI截图匹配，成功率>80%；替代人工截图匹配，节省1183人天。"},
        {"header": "智能问答助手", "body": "接入LLM解决3476次日常咨询，实现一站式知识查询。"},
        {"header": "缺陷去重", "body": "利用NLP技术识别重复缺陷，拦截2022次重复提单，大幅降低无效Bug流转。"}
    ])

    # Slide 7: Quick App Gov Service
    create_slide(prs, "快应用国家政务服务平台", [
        {"header": "高考/高频服务保障", "body": "支撑11省高考查分，31省录取查询，实现100%零故障。"},
        {"header": "核心业务指标", "body": "社保/医保服务 PV 90w+，UV 180w+；服务稳定性达99.99%。"},
        {"header": "技术优化", "body": "完善接口自动化监控；修复Nginx配置与系统时间同步导致的鉴权偶发失败问题。"}
    ])

    # Slide 8: Problems & Suggestions
    create_slide(prs, "问题与建议", [
        {"header": "技术痛点", "body": "系统间时间不同步导致OAuth Code偶发过期；部分旧系统接口缺乏监控覆盖。"},
        {"header": "改进建议", "body": "1. 运维协同优化全链路Nginx配置与NTP时间同步。\n2. 加强核心业务接口的自动化监控覆盖率，实现秒级告警。"},
        {"header": "流程优化", "body": "建议推广AI缺陷去重工具至更多项目组，减少重复劳动。"}
    ])

    # Slide 9: New Year Plan
    create_slide(prs, "2026年工作规划", [
        {"header": "快应用业务", "body": "落地医保移动支付功能；优化证照中心体验（卡面展示、交互升级）。"},
        {"header": "MEAT平台深化", "body": "深化AI在Code Review、需求分析中的应用；VLA支持更多日志格式解析。"},
        {"header": "团队建设", "body": "提升技术分享频率，沉淀更多通用前端组件；探索鸿蒙原生应用开发技术。"}
    ])

    # Slide 10: End
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "谢谢观看"
    slide.placeholders[1].text = "2025年年终总结"

    output_path = os.path.join(os.getcwd(), "2025_Year_End_Summary.pptx")
    prs.save(output_path)
    print(f"PPTX generated at: {output_path}")

if __name__ == "__main__":
    main()
